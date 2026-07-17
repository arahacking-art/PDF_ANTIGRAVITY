from fastapi import FastAPI, UploadFile, File, Form, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse, JSONResponse
from contextlib import asynccontextmanager
import pikepdf
import pymupdf  # PyMuPDF (fitz)
import io
import asyncio
import zipfile
import pdfplumber
import openpyxl
from pdf2docx import Converter
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import tempfile
import os
import subprocess
import base64
import json
from concurrent.futures import ProcessPoolExecutor

# ─── Semáforos de Concurrencia ────────────────────────────────────────────────
_ocr_semaphore: asyncio.Semaphore | None = None
_libreoffice_semaphore: asyncio.Semaphore | None = None
_convert_semaphore: asyncio.Semaphore | None = None

@asynccontextmanager
async def lifespan(app: FastAPI):
    global _ocr_semaphore, _libreoffice_semaphore, _convert_semaphore
    _ocr_semaphore        = asyncio.Semaphore(2)  # Máx 2 OCR simultáneos
    _libreoffice_semaphore = asyncio.Semaphore(2)  # Máx 2 conversiones LibreOffice
    _convert_semaphore    = asyncio.Semaphore(4)  # Máx 4 conversiones PDF simultáneas
    yield  # La app corre aquí
    # (cleanup si fuera necesario)

app = FastAPI(
    title="PDF Antigravity API",
    description="Suite empresarial de manipulación de PDFs. Procesos pesados asíncronos con semáforos de concurrencia.",
    version="3.1.0",
    lifespan=lifespan,
)

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
    expose_headers=["X-Original-Size", "X-Compressed-Size", "X-Reduction-Percent", "X-Result-Size"],
)



@app.get("/")
def read_root():
    return {"message": "Welcome to PDF Antigravity API v2.0"}


# ─── SEGURIDAD ────────────────────────────────────────────────────────────────

@app.post("/api/protect")
async def protect_pdf(file: UploadFile = File(...), password: str = Form(...)):
    try:
        contents = await file.read()
        pdf_in = io.BytesIO(contents)
        pdf_out = io.BytesIO()
        with pikepdf.open(pdf_in) as pdf:
            pdf.save(pdf_out, encryption=pikepdf.Encryption(owner=password, user=password, R=6))
        pdf_out.seek(0)
        return StreamingResponse(pdf_out, media_type="application/pdf",
            headers={"Content-Disposition": f'attachment; filename="protegido_{file.filename}"'})
    except pikepdf.PasswordError:
        raise HTTPException(400, "El PDF ya está protegido o está dañado.")
    except Exception as e:
        raise HTTPException(500, str(e))


@app.post("/api/unlock")
async def unlock_pdf(file: UploadFile = File(...), password: str = Form(...)):
    try:
        contents = await file.read()
        pdf_in = io.BytesIO(contents)
        pdf_out = io.BytesIO()
        with pikepdf.open(pdf_in, password=password) as pdf:
            pdf.save(pdf_out)
        pdf_out.seek(0)
        return StreamingResponse(pdf_out, media_type="application/pdf",
            headers={"Content-Disposition": f'attachment; filename="desbloqueado_{file.filename}"'})
    except pikepdf.PasswordError:
        raise HTTPException(401, "Contraseña incorrecta. Verifica e intenta de nuevo.")
    except Exception as e:
        raise HTTPException(500, str(e))


# ─── COMPRIMIR ────────────────────────────────────────────────────────────────

@app.post("/api/compress")
async def compress_pdf(file: UploadFile = File(...), mode: str = Form("standard")):
    """
    Comprime el PDF con múltiples estrategias. Procesado asíncronamente para no bloquear.
    """
    try:
        contents = await file.read()
        original_size = len(contents)

        def process_compression():
            best_result = io.BytesIO(contents)
            best_size = original_size

            # ── Estrategia 1: PyMuPDF garbage=4 + deflate completo ─────────────────
            try:
                doc = pymupdf.open(stream=contents, filetype="pdf")
                try:
                    doc.scrub()
                except Exception:
                    pass
                attempt1 = io.BytesIO()
                doc.save(attempt1, garbage=4, deflate=True, deflate_images=True, deflate_fonts=True, clean=True)
                doc.close()
                size1 = attempt1.getbuffer().nbytes
                if size1 < best_size:
                    best_size = size1
                    best_result = attempt1
            except Exception:
                pass

            # ── Estrategia 2: pikepdf con recompresión flate ────────────────────────
            try:
                attempt2 = io.BytesIO()
                with pikepdf.open(io.BytesIO(contents)) as pdf:
                    pdf.save(attempt2, compress_streams=True, object_stream_mode=pikepdf.ObjectStreamMode.generate, recompress_flate=True, preserve_pdfa=False)
                size2 = attempt2.getbuffer().nbytes
                if size2 < best_size:
                    best_size = size2
                    best_result = attempt2
            except Exception:
                pass

            # ── Estrategia 3 (solo modo "aggressive"): Resamplear imágenes ──────────
            if mode == "aggressive":
                try:
                    doc = pymupdf.open(stream=contents, filetype="pdf")
                    for page in doc:
                        for img_info in page.get_images(full=True):
                            xref = img_info[0]
                            try:
                                base_img = doc.extract_image(xref)
                                img_bytes = base_img["image"]
                                if len(img_bytes) < 50_000:
                                    continue
                                pil_img = Image.open(io.BytesIO(img_bytes))
                                if max(pil_img.width, pil_img.height) > 1200:
                                    pil_img.thumbnail((1200, 1200), Image.LANCZOS)
                                out_buf = io.BytesIO()
                                pil_img.convert("RGB").save(out_buf, format="JPEG", quality=75, optimize=True)
                                doc.update_stream(xref, out_buf.getvalue())
                            except Exception:
                                continue
                    attempt3 = io.BytesIO()
                    doc.save(attempt3, garbage=4, deflate=True, deflate_images=True, deflate_fonts=True, clean=True)
                    doc.close()
                    size3 = attempt3.getbuffer().nbytes
                    if size3 < best_size:
                        best_size = size3
                        best_result = attempt3
                except Exception:
                    pass
            return best_result, best_size

        async with _convert_semaphore:
            best_result, best_size = await asyncio.to_thread(process_compression)
            
        best_result.seek(0)
        reduction_pct = max(0.0, (original_size - best_size) / original_size * 100)

        return StreamingResponse(
            best_result,
            media_type="application/pdf",
            headers={
                "Content-Disposition": f'attachment; filename="comprimido_{file.filename}"',
                "X-Original-Size": str(original_size),
                "X-Compressed-Size": str(best_size),
                "X-Reduction-Percent": f"{reduction_pct:.1f}",
                "Access-Control-Expose-Headers": "X-Original-Size,X-Compressed-Size,X-Reduction-Percent",
            },
        )
    except Exception as e:
        raise HTTPException(500, f"Error al comprimir: {str(e)}")


# ─── CONVERSIONES ─────────────────────────────────────────────────────────────

@app.post("/api/to-word")
async def pdf_to_word(file: UploadFile = File(...)):
    """Convierte PDF a DOCX usando pdf2docx."""
    try:
        contents = await file.read()
        
        def process_word():
            with tempfile.TemporaryDirectory() as tmpdir:
                pdf_path = os.path.join(tmpdir, "input.pdf")
                docx_path = os.path.join(tmpdir, "output.docx")
                with open(pdf_path, "wb") as f:
                    f.write(contents)
                cv = Converter(pdf_path)
                cv.convert(docx_path)
                cv.close()
                with open(docx_path, "rb") as f:
                    return f.read()

        async with _convert_semaphore:
            docx_bytes = await asyncio.to_thread(process_word)

        docx_name = file.filename.replace(".pdf", ".docx") if file.filename else "output.docx"
        return StreamingResponse(
            io.BytesIO(docx_bytes),
            media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            headers={"Content-Disposition": f'attachment; filename="{docx_name}"'},
        )
    except Exception as e:
        raise HTTPException(500, f"Error al convertir a Word: {str(e)}")


@app.post("/api/to-pptx")
async def pdf_to_pptx(file: UploadFile = File(...)):
    """Convierte cada página del PDF en una diapositiva con imagen (PyMuPDF + python-pptx)."""
    try:
        contents = await file.read()
        
        def process_pptx():
            doc = pymupdf.open(stream=contents, filetype="pdf")
            prs = Presentation()
            prs.slide_width = Inches(13.33)
            prs.slide_height = Inches(7.5)
            blank_layout = prs.slide_layouts[6]

            with tempfile.TemporaryDirectory() as tmpdir:
                for i, page in enumerate(doc):
                    mat = pymupdf.Matrix(2.0, 2.0)
                    pix = page.get_pixmap(matrix=mat)
                    img_path = os.path.join(tmpdir, f"page_{i}.png")
                    pix.save(img_path)
                    slide = prs.slides.add_slide(blank_layout)
                    slide.shapes.add_picture(img_path, 0, 0, prs.slide_width, prs.slide_height)

                pptx_path = os.path.join(tmpdir, "output.pptx")
                prs.save(pptx_path)
                with open(pptx_path, "rb") as f:
                    return f.read()

        async with _convert_semaphore:
            pptx_bytes = await asyncio.to_thread(process_pptx)

        pptx_name = file.filename.replace(".pdf", ".pptx") if file.filename else "output.pptx"
        return StreamingResponse(
            io.BytesIO(pptx_bytes),
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": f'attachment; filename="{pptx_name}"'},
        )
    except Exception as e:
        raise HTTPException(500, f"Error al convertir a PowerPoint: {str(e)}")


@app.post("/api/to-excel")
async def pdf_to_excel(file: UploadFile = File(...)):
    """Extrae tablas del PDF y las exporta a XLSX usando pdfplumber + openpyxl."""
    try:
        contents = await file.read()
        
        def process_excel():
            wb = openpyxl.Workbook()
            wb.remove(wb.active)
            tables_found = 0
            with pdfplumber.open(io.BytesIO(contents)) as pdf:
                for page_num, page in enumerate(pdf.pages, start=1):
                    tables = page.extract_tables()
                    for table_idx, table in enumerate(tables):
                        tables_found += 1
                        sheet_name = f"Pag{page_num}_T{table_idx + 1}"
                        ws = wb.create_sheet(title=sheet_name)
                        for row in table:
                            ws.append([cell if cell is not None else "" for cell in row])
            if tables_found == 0:
                ws = wb.create_sheet(title="Texto_Extraido")
                with pdfplumber.open(io.BytesIO(contents)) as pdf:
                    for page_num, page in enumerate(pdf.pages, start=1):
                        text = page.extract_text() or ""
                        ws.append([f"--- Página {page_num} ---"])
                        for line in text.split("\n"):
                            ws.append([line])
                        ws.append([])
            xlsx_out = io.BytesIO()
            wb.save(xlsx_out)
            xlsx_out.seek(0)
            return xlsx_out
            
        async with _convert_semaphore:
            xlsx_out = await asyncio.to_thread(process_excel)

        xlsx_name = file.filename.replace(".pdf", ".xlsx") if file.filename else "output.xlsx"
        return StreamingResponse(
            xlsx_out,
            media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            headers={"Content-Disposition": f'attachment; filename="{xlsx_name}"'},
        )
    except Exception as e:
        raise HTTPException(500, f"Error al convertir a Excel: {str(e)}")


@app.post("/api/to-jpg")
async def pdf_to_jpg(file: UploadFile = File(...), dpi: int = Form(150)):
    """
    Convierte páginas del PDF a imágenes JPG.
    1 página → devuelve JPG directo.
    Múltiples páginas → devuelve un ZIP con todos los JPGs.
    Usa asyncio.to_thread para no bloquear el event loop.
    """
    try:
        contents = await file.read()
        base_name = file.filename.replace(".pdf", "") if file.filename else "pdf"

        def process_jpg():
            doc = pymupdf.open(stream=contents, filetype="pdf")
            zoom = dpi / 72.0
            mat = pymupdf.Matrix(zoom, zoom)
            images: list[tuple[str, bytes]] = []
            for i, page in enumerate(doc):
                pix = page.get_pixmap(matrix=mat)
                img_bytes = pix.tobytes("jpeg")
                images.append((f"pagina_{i + 1:03d}.jpg", img_bytes))
            doc.close()
            return images

        async with _convert_semaphore:
            images = await asyncio.to_thread(process_jpg)

        if len(images) == 1:
            return StreamingResponse(
                io.BytesIO(images[0][1]),
                media_type="image/jpeg",
                headers={"Content-Disposition": f'attachment; filename="{base_name}.jpg"'},
            )

        # Múltiples páginas → ZIP
        zip_out = io.BytesIO()
        with zipfile.ZipFile(zip_out, "w", zipfile.ZIP_DEFLATED) as zf:
            for name, data in images:
                zf.writestr(name, data)
        zip_out.seek(0)
        return StreamingResponse(
            zip_out,
            media_type="application/zip",
            headers={"Content-Disposition": f'attachment; filename="{base_name}_imagenes.zip"'},
        )
    except Exception as e:
        raise HTTPException(500, f"Error al convertir a JPG: {str(e)}")



# ─── CENSURAR (REDACTAR) ──────────────────────────────────────────────────────

@app.post("/api/redact")
async def redact_pdf(file: UploadFile = File(...), terms: str = Form(...)):
    """
    Busca y censura (cubre con negro) los términos indicados en el PDF.
    `terms` es una lista separada por comas o saltos de línea.
    """
    try:
        contents = await file.read()
        doc = pymupdf.open(stream=contents, filetype="pdf")

        search_terms = [t.strip() for t in terms.replace("\n", ",").split(",") if t.strip()]
        if not search_terms:
            raise HTTPException(400, "Debes indicar al menos un término para censurar.")

        redactions_applied = 0
        for page in doc:
            for term in search_terms:
                hits = page.search_for(term, quads=True)
                for hit in hits:
                    page.add_redact_annot(hit, fill=(0, 0, 0))
                    redactions_applied += 1
            if redactions_applied > 0:
                page.apply_redactions()

        pdf_out = io.BytesIO()
        doc.save(pdf_out)
        doc.close()
        pdf_out.seek(0)
        return StreamingResponse(pdf_out, media_type="application/pdf",
            headers={"Content-Disposition": f'attachment; filename="censurado_{file.filename}"'})
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500, f"Error al censurar: {str(e)}")


# ─── RENDERIZADO VISUAL ───────────────────────────────────────────────────────

@app.post("/api/render-pages")
async def render_pages(
    file: UploadFile = File(...),
    scale: float = Form(1.5),
):
    """
    Renderiza cada página del PDF como PNG en base64.
    Devuelve dimensiones originales (PDF points) y de renderizado (px).
    """
    try:
        contents = await file.read()
        doc = pymupdf.open(stream=contents, filetype="pdf")
        mat = pymupdf.Matrix(scale, scale)
        pages_data = []
        for i, page in enumerate(doc):
            pix = page.get_pixmap(matrix=mat)
            img_b64 = base64.b64encode(pix.tobytes("png")).decode()
            pages_data.append({
                "index": i,
                "width": pix.width,
                "height": pix.height,
                "originalWidth": page.rect.width,
                "originalHeight": page.rect.height,
                "data": img_b64,
            })
        doc.close()
        return JSONResponse({"pages": pages_data, "pageCount": len(pages_data)})
    except Exception as e:
        raise HTTPException(500, f"Error al renderizar páginas: {str(e)}")


@app.post("/api/redact-visual")
async def redact_visual(
    file: UploadFile = File(...),
    areas: str = Form(...),  # JSON: [{page, x1r, y1r, x2r, y2r}] — ratios 0-1
):
    """
    Censura áreas en el PDF usando coordenadas relativas (ratios).
    x1r, y1r, x2r, y2r son proporción del ancho/alto de la página (0.0 a 1.0).
    El sistema de coordenadas de PyMuPDF tiene origen en la esquina superior izquierda.
    """
    try:
        contents = await file.read()
        redact_areas = json.loads(areas)
        doc = pymupdf.open(stream=contents, filetype="pdf")

        for area in redact_areas:
            page_idx = int(area["page"])
            if 0 <= page_idx < len(doc):
                page = doc[page_idx]
                pw = page.rect.width
                ph = page.rect.height
                x1 = float(area["x1r"]) * pw
                y1 = float(area["y1r"]) * ph
                x2 = float(area["x2r"]) * pw
                y2 = float(area["y2r"]) * ph
                page.add_redact_annot(pymupdf.Rect(x1, y1, x2, y2), fill=(0, 0, 0))
                page.apply_redactions()

        pdf_out = io.BytesIO()
        doc.save(pdf_out)
        doc.close()
        pdf_out.seek(0)
        return StreamingResponse(pdf_out, media_type="application/pdf",
            headers={"Content-Disposition": f'attachment; filename="censurado_{file.filename}"'})
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500, f"Error al censurar visualmente: {str(e)}")


@app.post("/api/sign-visual")
async def sign_visual(
    file: UploadFile = File(...),
    signature: UploadFile = File(...),
    position: str = Form(...),  # JSON: {page, x_ratio, y_ratio, w_ratio, h_ratio}
):
    """
    Incrusta la imagen de firma en la página y posición indicadas.
    Las coordenadas son ratios (0-1) del ancho/alto de la página.
    """
    try:
        pdf_bytes = await file.read()
        sig_bytes = await signature.read()
        pos = json.loads(position)

        doc = pymupdf.open(stream=pdf_bytes, filetype="pdf")
        page_idx = int(pos["page"])
        page = doc[page_idx]

        pw = page.rect.width
        ph = page.rect.height
        x = float(pos["x_ratio"]) * pw
        y = float(pos["y_ratio"]) * ph
        w = float(pos["w_ratio"]) * pw
        h = float(pos["h_ratio"]) * ph

        rect = pymupdf.Rect(x, y, x + w, y + h)
        page.insert_image(rect, stream=sig_bytes, keep_proportion=True)

        pdf_out = io.BytesIO()
        doc.save(pdf_out)
        doc.close()
        pdf_out.seek(0)
        return StreamingResponse(pdf_out, media_type="application/pdf",
            headers={"Content-Disposition": f'attachment; filename="firmado_{file.filename}"'})
    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(500, f"Error al firmar visualmente: {str(e)}")


# ─── OFFICE → PDF ─────────────────────────────────────────────────────────────

@app.post("/api/office-to-pdf")
async def office_to_pdf(file: UploadFile = File(...)):
    """
    Convierte archivos Office (DOCX, XLSX, PPTX, ODT, CSV, TXT) a PDF
    usando LibreOffice en modo headless.
    Soporta: .docx .doc .odt .xlsx .xls .ods .pptx .ppt .odp .csv .txt .rtf
    """
    ALLOWED_EXTENSIONS = {
        ".docx", ".doc", ".odt", ".rtf",   # Writer
        ".xlsx", ".xls", ".ods", ".csv",    # Calc
        ".pptx", ".ppt", ".odp",            # Impress
        ".txt",
    }

    filename = file.filename or "documento"
    ext = os.path.splitext(filename)[1].lower()

    if ext not in ALLOWED_EXTENSIONS:
        raise HTTPException(
            400,
            f"Formato no soportado: '{ext}'. "
            f"Formatos aceptados: {', '.join(sorted(ALLOWED_EXTENSIONS))}"
        )

    try:
        contents = await file.read()

        def process_office():
            with tempfile.TemporaryDirectory() as tmpdir:
                input_path = os.path.join(tmpdir, filename)
                with open(input_path, "wb") as f:
                    f.write(contents)
                
                try:
                    result = subprocess.run(
                        ["libreoffice", "--headless", "--norestore", "--nofirststartwizard", "--convert-to", "pdf", "--outdir", tmpdir, input_path],
                        capture_output=True, text=True, timeout=120
                    )
                except subprocess.TimeoutExpired:
                    raise RuntimeError("TIMEOUT")
                
                if result.returncode != 0:
                    raise RuntimeError(f"LibreOffice falló al convertir: {result.stderr.strip() or 'Error desconocido'}")
                
                pdf_name = os.path.splitext(filename)[0] + ".pdf"
                pdf_path = os.path.join(tmpdir, pdf_name)
                if not os.path.exists(pdf_path):
                    raise RuntimeError("LibreOffice no generó el PDF. Verifica que el archivo no esté dañado.")
                
                with open(pdf_path, "rb") as f:
                    return f.read(), pdf_name

        async with _libreoffice_semaphore:
            pdf_bytes, pdf_name = await asyncio.to_thread(process_office)

        return StreamingResponse(
            io.BytesIO(pdf_bytes),
            media_type="application/pdf",
            headers={"Content-Disposition": f'attachment; filename="{pdf_name}"'},
        )
    except Exception as e:
        if "TIMEOUT" in str(e):
            raise HTTPException(504, "La conversión tardó demasiado (límite: 120 segundos). El archivo puede ser muy complejo.")
        raise HTTPException(500, f"Error al convertir: {str(e)}")


# ─── OCR ──────────────────────────────────────────────────────────────────────

@app.post("/api/ocr")
async def ocr_pdf(
    file: UploadFile = File(...),
    language: str = Form("spa+eng"),
    deskew: bool = Form(True),
):
    """Aplica OCR asíncronamente con protección de concurrencia."""
    import ocrmypdf
    try:
        contents = await file.read()

        def process_ocr():
            with tempfile.TemporaryDirectory() as tmpdir:
                input_path = os.path.join(tmpdir, "input.pdf")
                output_path = os.path.join(tmpdir, "output_ocr.pdf")
                with open(input_path, "wb") as f:
                    f.write(contents)
                ocrmypdf.ocr(
                    input_path, output_path, language=language, deskew=deskew, skip_text=True, optimize=1, oversample=300, progress_bar=False
                )
                with open(output_path, "rb") as f:
                    return f.read()

        async with _ocr_semaphore:
            result_bytes = await asyncio.to_thread(process_ocr)

        filename_out = file.filename.replace(".pdf", "_OCR.pdf") if file.filename else "documento_OCR.pdf"
        return StreamingResponse(
            io.BytesIO(result_bytes),
            media_type="application/pdf",
            headers={
                "Content-Disposition": f'attachment; filename="{filename_out}"',
                "X-Original-Size": str(len(contents)),
                "X-Result-Size": str(len(result_bytes)),
                "Access-Control-Expose-Headers": "X-Original-Size,X-Result-Size",
            },
        )
    except ocrmypdf.exceptions.PriorOcrFoundError:
        raise HTTPException(409, "Este PDF ya contiene texto seleccionable. Usa 'Forzar re-OCR'.")
    except ocrmypdf.exceptions.EncryptedPdfError:
        raise HTTPException(400, "El PDF está protegido con contraseña. Quita la contraseña primero.")
    except Exception as e:
        raise HTTPException(500, f"Error durante el OCR: {str(e)}")


@app.post("/api/ocr-force")
async def ocr_pdf_force(
    file: UploadFile = File(...),
    language: str = Form("spa+eng"),
    deskew: bool = Form(True),
):
    """Igual que /api/ocr pero fuerza el re-procesamiento."""
    import ocrmypdf
    try:
        contents = await file.read()

        def process_ocr_force():
            with tempfile.TemporaryDirectory() as tmpdir:
                input_path = os.path.join(tmpdir, "input.pdf")
                output_path = os.path.join(tmpdir, "output_ocr.pdf")
                with open(input_path, "wb") as f:
                    f.write(contents)
                ocrmypdf.ocr(
                    input_path, output_path, language=language, deskew=deskew, force_ocr=True, optimize=1, progress_bar=False
                )
                with open(output_path, "rb") as f:
                    return f.read()

        async with _ocr_semaphore:
            result_bytes = await asyncio.to_thread(process_ocr_force)

        filename_out = file.filename.replace(".pdf", "_OCR.pdf") if file.filename else "documento_OCR.pdf"
        return StreamingResponse(
            io.BytesIO(result_bytes), media_type="application/pdf",
            headers={"Content-Disposition": f'attachment; filename="{filename_out}"'},
        )
    except Exception as e:
        raise HTTPException(500, f"Error durante el OCR forzado: {str(e)}")


